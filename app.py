import os
import sys
import warnings
warnings.filterwarnings('ignore')

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding='utf-8')
        sys.stderr.reconfigure(encoding='utf-8')
    except Exception:
        pass

os.environ['HF_HUB_DISABLE_SYMLINKS_WARNING'] = '1'
os.environ['TOKENIZERS_PARALLELISM'] = 'false'

from pptx import Presentation
from flask import Flask, render_template, request, jsonify, session, redirect, url_for, Response
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
from huggingface_hub import HfApi, InferenceClient
from pymongo import MongoClient
import PyPDF2
import docx
import pandas as pd
from dotenv import load_dotenv
from functools import wraps
import json
from datetime import timedelta, datetime
import traceback
import time
import re
import uuid
import requests
import sqlite3

load_dotenv()

# ========== HUGGING FACE SECRETS SUPPORT ==========
def load_hf_secrets():
    secrets_dir = '/run/secrets'
    if os.path.exists(secrets_dir):
        for secret_name in os.listdir(secrets_dir):
            secret_path = os.path.join(secrets_dir, secret_name)
            try:
                with open(secret_path, 'r') as f:
                    os.environ[secret_name] = f.read().strip()
                print(f"✅ Loaded secret: {secret_name}")
            except Exception as e:
                print(f"⚠️ Could not load {secret_name}: {e}")

    for secret_name in ['GEMINI_API_KEY', 'GROQ_API_KEY', 'SECRET_KEY', 'ADMIN_USERNAME', 'ADMIN_PASSWORD',
                         'QDRANT_URL', 'QDRANT_API_KEY', 'HF_TOKEN', 'HF_DATASET', 'MONGO_URI']:
        if not os.environ.get(secret_name):
            paths = [f'/etc/secrets/{secret_name}', f'/secrets/{secret_name}', f'/run/secrets/{secret_name}']
            for path in paths:
                if os.path.exists(path):
                    try:
                        with open(path, 'r') as f:
                            os.environ[secret_name] = f.read().strip()
                        print(f"✅ Loaded secret: {secret_name}")
                        break
                    except:
                        pass

load_hf_secrets()

gemini_key = os.environ.get('GEMINI_API_KEY', 'NOT SET')
groq_key = os.environ.get('GROQ_API_KEY', 'NOT SET')
mongo_uri = os.environ.get('MONGO_URI', 'NOT SET')
print(f"🔑 GEMINI_API_KEY: {'SET' if gemini_key != 'NOT SET' else 'NOT SET'}")
print(f"🔑 GROQ_API_KEY: {'SET' if groq_key != 'NOT SET' else 'NOT SET'}")
print(f"🔑 MONGO_URI: {'SET' if mongo_uri != 'NOT SET' else 'NOT SET'}")
print(f"🔑 QDRANT_URL: {'SET' if os.environ.get('QDRANT_URL') else 'NOT SET'}")

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = 'uploads'
app.secret_key = os.getenv('SECRET_KEY', 'your-secret-key-change-this')
app.permanent_session_lifetime = timedelta(hours=24)

ADMIN_USERNAME = os.getenv('ADMIN_USERNAME', 'admin')
ADMIN_PASSWORD_HASH = generate_password_hash(os.getenv('ADMIN_PASSWORD', 'Admin@123'))

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# ========== MONGODB FOR STUDENT ACCOUNTS & CHAT SESSIONS ==========
mongo_client = None
users_collection = None
sessions_collection = None
mongo_available = False

try:
    if mongo_uri != 'NOT SET':
        mongo_client = MongoClient(mongo_uri, serverSelectionTimeoutMS=5000)
        mongo_client.admin.command('ping')
        db_mongo = mongo_client['chatbot_db']
        users_collection = db_mongo['users']
        sessions_collection = db_mongo['chat_sessions']
        users_collection.create_index('email', unique=True)
        mongo_available = True
        print("✅ MongoDB connected (Users & Chat Sessions)")
    else:
        print("⚠️ MONGO_URI not set - student accounts disabled")
except Exception as e:
    print(f"⚠️ MongoDB connection failed: {str(e)[:100]}")
    print("⚠️ Student accounts disabled - using SQLite fallback for chat")

def safe_mongo_count(collection):
    if collection is None:
        return 0
    try:
        return collection.count_documents({})
    except Exception:
        return 0

# ========== PERSISTENT CONVERSATION MEMORY (SQLite fallback) ==========
DB_PATH = os.getenv('CHAT_DB_PATH', 'chat_history.db')
MAX_HISTORY = 20

def get_db_connection():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False, timeout=10)
    conn.execute('PRAGMA journal_mode=WAL')
    return conn

def init_db():
    conn = get_db_connection()
    conn.execute('''
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            timestamp TEXT NOT NULL
        )
    ''')
    conn.execute('CREATE INDEX IF NOT EXISTS idx_session_id ON messages(session_id)')
    conn.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            email TEXT NOT NULL UNIQUE,
            password_hash TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
    ''')
    conn.commit()
    conn.close()

def save_message_sqlite(session_id, role, content):
    conn = get_db_connection()
    try:
        conn.execute(
            'INSERT INTO messages (session_id, role, content, timestamp) VALUES (?, ?, ?, ?)',
            (session_id, role, content, datetime.utcnow().isoformat())
        )
        conn.commit()
    finally:
        conn.close()

def get_session_history_sqlite(session_id, limit=MAX_HISTORY):
    conn = get_db_connection()
    try:
        cur = conn.execute(
            'SELECT role, content FROM messages WHERE session_id = ? ORDER BY id DESC LIMIT ?',
            (session_id, limit)
        )
        rows = cur.fetchall()
    finally:
        conn.close()
    return [{'role': r, 'content': c} for r, c in reversed(rows)]

def trim_session_history_sqlite(session_id, keep=MAX_HISTORY):
    conn = get_db_connection()
    try:
        conn.execute('''
            DELETE FROM messages
            WHERE session_id = ? AND id NOT IN (
                SELECT id FROM messages WHERE session_id = ? ORDER BY id DESC LIMIT ?
            )
        ''', (session_id, session_id, keep))
        conn.commit()
    finally:
        conn.close()

def clear_session_history_sqlite(session_id):
    conn = get_db_connection()
    try:
        conn.execute('DELETE FROM messages WHERE session_id = ?', (session_id,))
        conn.commit()
    finally:
        conn.close()

def save_message(session_id, role, content, sources=None):
    if mongo_available and sessions_collection is not None:
        try:
            sessions_collection.update_one(
                {'session_id': session_id},
                {'$push': {'messages': {
                    'role': role, 'content': content,
                    'sources': sources or [], 'timestamp': datetime.utcnow()
                }}},
                upsert=True
            )
            return
        except Exception as e:
            print(f"MongoDB save error: {e}")
    save_message_sqlite(session_id, role, content)

def get_session_history(session_id, limit=MAX_HISTORY):
    if mongo_available and sessions_collection is not None:
        try:
            doc = sessions_collection.find_one({'session_id': session_id})
            if doc and 'messages' in doc:
                messages = doc['messages'][-limit:]
                return [{'role': m['role'], 'content': m['content']} for m in messages]
        except Exception as e:
            print(f"MongoDB get error: {e}")
    return get_session_history_sqlite(session_id, limit)

def trim_session_history(session_id, keep=MAX_HISTORY):
    if mongo_available and sessions_collection is not None:
        try:
            doc = sessions_collection.find_one({'session_id': session_id})
            if doc and 'messages' in doc and len(doc['messages']) > keep:
                sessions_collection.update_one(
                    {'session_id': session_id},
                    {'$push': {'messages': {'$each': [], '$slice': -keep}}}
                )
            return
        except Exception as e:
            print(f"MongoDB trim error: {e}")
    trim_session_history_sqlite(session_id, keep)

def clear_session_history(session_id):
    if mongo_available and sessions_collection is not None:
        try:
            sessions_collection.delete_one({'session_id': session_id})
            return
        except Exception as e:
            print(f"MongoDB clear error: {e}")
    clear_session_history_sqlite(session_id)

init_db()

print("\n" + "="*60)
print("🔄 INITIALIZING...")
print("="*60)

class HFEmbeddingModel:
    def __init__(self, model_name="sentence-transformers/all-MiniLM-L6-v2", token=None):
        self.model_name = model_name
        self.token = token
        self.local_model = None
        self.client = None
        
        if SentenceTransformer is not None:
            try:
                self.local_model = SentenceTransformer(self.model_name)
                print(f"✅ Loaded local SentenceTransformer model: {self.model_name}")
            except Exception as e:
                print(f"⚠️ Could not load local SentenceTransformer: {e}. Falling back to Hugging Face Inference API.")
        else:
            print("⚠️ sentence-transformers package not installed. Falling back to Hugging Face Inference API.")
            
        if self.local_model is None:
            self.client = InferenceClient(token=self.token)
            print("✅ Initialized Hugging Face InferenceClient")

    def encode(self, sentences, **kwargs):
        if self.local_model:
            return self.local_model.encode(sentences, **kwargs)
        if not self.client:
            self.client = InferenceClient(token=self.token)
        try:
            import numpy as np
            if isinstance(sentences, str):
                res = self.client.feature_extraction(text=sentences, model=self.model_name)
                val = res if isinstance(res, np.ndarray) else np.array(res)
                if np.all(val == 0):
                    raise ValueError("Inference API returned all zeros")
                return val
            else:
                results = []
                for s in sentences:
                    res = self.client.feature_extraction(text=s, model=self.model_name)
                    val = res if isinstance(res, np.ndarray) else np.array(res)
                    if np.all(val == 0):
                        raise ValueError("Inference API returned all zeros")
                    results.append(val)
                return np.stack(results)
        except Exception as e:
            print(f"❌ Error during HF Inference API feature extraction: {e}")
            raise RuntimeError(f"Embedding generation failed: {e}")

try:
    hf_token = os.getenv('HF_TOKEN')
    embedding_model = HFEmbeddingModel(token=hf_token)
    print("✅ Hybrid embedding model loaded successfully")
except Exception as e:
    embedding_model = None
    print(f"❌ Embedding failed: {e}")

qdrant_client = None
try:
    qdrant_url = os.getenv('QDRANT_URL')
    qdrant_api_key = os.getenv('QDRANT_API_KEY')
    if qdrant_url:
        qdrant_client = QdrantClient(url=qdrant_url, api_key=qdrant_api_key if qdrant_api_key else None)
        print("✅ Qdrant connected to Cloud/Docker instance")
    else:
        db_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "qdrant_db")
        qdrant_client = QdrantClient(path=db_dir)
        print(f"📁 Qdrant local storage initialized at: {db_dir}")

    if qdrant_client:
        collections = qdrant_client.get_collections().collections
        collection_names = [c.name for c in collections]
        if "university_notes" not in collection_names:
            qdrant_client.create_collection(
                collection_name="university_notes",
                vectors_config=VectorParams(size=384, distance=Distance.COSINE)
            )
            print("✅ Qdrant collection created")
        else:
            info = qdrant_client.get_collection("university_notes")
            print(f"✅ Qdrant connected ({info.points_count} documents)")
            
        # Ensure payload index for 'source' is created
        try:
            qdrant_client.create_payload_index(
                collection_name="university_notes",
                field_name="source",
                field_schema="keyword"
            )
            print("✅ Created Qdrant payload index for 'source'")
        except Exception as e:
            print(f"⚠️ Could not create Qdrant payload index: {e}")
except Exception as e:
    print(f"❌ Qdrant connection failed: {e}")
    qdrant_client = None

hf_api = None
hf_dataset = os.getenv('HF_DATASET', '')
try:
    hf_token = os.getenv('HF_TOKEN')
    if hf_token and hf_dataset:
        hf_api = HfApi(token=hf_token)
        print(f"✅ HF Dataset ready: {hf_dataset}")
    else:
        print("⚠️ HF Dataset credentials not found")
except Exception as e:
    print(f"⚠️ HF Dataset setup failed: {e}")

gemini_api_key = os.getenv('GEMINI_API_KEY')
gemini_connected = bool(gemini_api_key)
groq_api_key = os.getenv('GROQ_API_KEY')
groq_connected = bool(groq_api_key)

if gemini_connected:
    print("✅ Gemini API key found")
if groq_connected:
    print("✅ Groq API key found")
if not gemini_connected and not groq_connected:
    print("❌ No AI API keys found!")

print("="*60 + "\n")

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'admin_logged_in' not in session:
            return redirect(url_for('admin_login_page'))
        return f(*args, **kwargs)
    return decorated_function

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'student_user' not in session:
            return redirect(url_for('login_page'))
        return f(*args, **kwargs)
    return decorated_function

# ========== FILE EXTRACTION ==========

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text and len(page_text.strip()) > 10:
                    text += page_text + "\n"
    except:
        pass
        
    # Fallback to pdfplumber which is highly robust for PowerPoint slide PDFs
    if len(text.strip()) < 100:
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                pdf_text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        pdf_text += page_text + "\n"
                if len(pdf_text.strip()) > len(text):
                    text = pdf_text
        except Exception as e:
            print(f"pdfplumber error: {e}")

    # Final fallback to OCR tesseract
    if len(text.strip()) < 100:
        try:
            from pdf2image import convert_from_path
            import pytesseract
            images = convert_from_path(file_path, dpi=150)
            ocr_text = ""
            for img in images:
                page_text = pytesseract.image_to_string(img)
                if page_text.strip():
                    ocr_text += page_text + "\n"
            if len(ocr_text) > len(text):
                text = ocr_text
        except:
            pass
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text += row_text + "\n"
            if slide_text.strip():
                text += f"\n--- Slide {slide_num} ---\n{slide_text}\n"
    except Exception as e:
        print(f"PPTX extraction error: {e}")
    return text

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except:
        return ""

def extract_text_from_txt(file_path):
    for enc in ['utf-8', 'latin-1', 'cp1252']:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except:
            continue
    return ""

def clean_text(text):
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'image\[\[.*?\]\]', '', text)
    text = re.sub(r'\b\d{1,2}/\d{1,2}/\d{2,4}\b', '', text)
    text = re.sub(r'<center>.*?</center>', '', text, flags=re.DOTALL)
    return text.strip()

def chunk_text(text, size=1000, overlap=150):
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current = ""
    for para in paragraphs:
        para = para.strip()
        if not para or len(para) < 20:
            continue
        if len(current) + len(para) < size:
            current += para + "\n\n"
        else:
            if current.strip():
                chunks.append(current.strip())
            current = para + "\n\n"
    if current.strip():
        chunks.append(current.strip())
    return chunks if chunks else [text[:size]]

# ========== API CLIENTS ==========

GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODELS = ["llama-3.1-8b-instant", "llama-3.3-70b-versatile", "mixtral-8x7b-32768"]
GEMINI_MODELS = ["gemini-2.0-flash", "gemini-2.0-flash-lite"]

def groq_chat_completion(messages, model="llama-3.1-8b-instant", max_tokens=150, temperature=0.7):
    headers = {"Authorization": f"Bearer {groq_api_key}", "Content-Type": "application/json"}
    payload = {"model": model, "messages": messages, "max_tokens": max_tokens, "temperature": temperature}
    response = requests.post(GROQ_URL, json=payload, headers=headers, timeout=30)
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    raise Exception(f"Groq API error: {response.status_code}")

def groq_chat_completion_stream(messages, model="llama-3.1-8b-instant", max_tokens=150, temperature=0.7):
    headers = {"Authorization": f"Bearer {groq_api_key}", "Content-Type": "application/json"}
    payload = {"model": model, "messages": messages, "max_tokens": max_tokens, "temperature": temperature, "stream": True}
    with requests.post(GROQ_URL, json=payload, headers=headers, timeout=30, stream=True) as response:
        if response.status_code != 200:
            raise Exception(f"Groq API error: {response.status_code}")
        for raw_line in response.iter_lines():
            if not raw_line: continue
            line = raw_line.decode('utf-8')
            if not line.startswith('data: '): continue
            payload_str = line[len('data: '):].strip()
            if payload_str == '[DONE]': break
            try:
                chunk = json.loads(payload_str)
                delta = chunk.get('choices', [{}])[0].get('delta', {}).get('content', '')
                if delta: yield delta
            except (json.JSONDecodeError, IndexError, KeyError): continue

def build_gemini_contents(history_messages, user_prompt):
    contents = []
    for msg in history_messages:
        role = "user" if msg['role'] == 'user' else "model"
        contents.append({"role": role, "parts": [{"text": msg['content']}]})
    contents.append({"role": "user", "parts": [{"text": user_prompt}]})
    return contents

def gemini_chat_completion(system_prompt, contents, model="gemini-2.0-flash", max_tokens=150, temperature=0.7):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={gemini_key}"
    headers = {"Content-Type": "application/json"}
    payload = {"contents": contents, "generationConfig": {"temperature": temperature, "maxOutputTokens": max_tokens}}
    if system_prompt:
        payload["systemInstruction"] = {"parts": [{"text": system_prompt}]}
    response = requests.post(url, json=payload, headers=headers, timeout=30)
    if response.status_code == 200:
        data = response.json()
        try:
            return data["candidates"][0]["content"]["parts"][0]["text"]
        except (KeyError, IndexError):
            raise Exception("Unexpected Gemini API response structure")
    raise Exception(f"Gemini API error: {response.status_code}")

def gemini_chat_completion_stream(system_prompt, contents, model="gemini-2.0-flash", max_tokens=150, temperature=0.7):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:streamGenerateContent?alt=sse&key={gemini_key}"
    headers = {"Content-Type": "application/json"}
    payload = {"contents": contents, "generationConfig": {"temperature": temperature, "maxOutputTokens": max_tokens}}
    if system_prompt:
        payload["systemInstruction"] = {"parts": [{"text": system_prompt}]}
    with requests.post(url, json=payload, headers=headers, timeout=30, stream=True) as response:
        if response.status_code != 200:
            raise Exception(f"Gemini API error: {response.status_code}")
        for raw_line in response.iter_lines():
            if not raw_line: continue
            line = raw_line.decode('utf-8')
            if not line.startswith('data: '): continue
            payload_str = line[len('data: '):].strip()
            try:
                chunk = json.loads(payload_str)
                delta = chunk["candidates"][0]["content"]["parts"][0]["text"]
                if delta: yield delta
            except (json.JSONDecodeError, KeyError, IndexError): continue

# ========== WEB SEARCH ==========

def search_duckduckgo(query):
    # Try using the library first
    try:
        from duckduckgo_search import DDGS
        results = []
        with DDGS() as ddgs:
            for r in list(ddgs.text(query, max_results=3)):
                results.append({"title": r.get("title", ""), "snippet": r.get("body", ""), "link": r.get("href", ""), "source": "DuckDuckGo"})
        if results:
            return results
    except Exception as e:
        print(f"DuckDuckGo API error: {e}")

    # Fallback: Scrape html.duckduckgo.com
    try:
        import urllib.parse
        from bs4 import BeautifulSoup
        url = "https://html.duckduckgo.com/html/"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        r = requests.get(url, params={"q": query}, headers=headers, timeout=8)
        if r.status_code == 200:
            soup = BeautifulSoup(r.text, 'html.parser')
            results = []
            for a in soup.find_all('a', class_='result__snippet')[:3]:
                parent = a.parent.parent
                title_el = parent.find('a', class_='result__url')
                snippet = a.get_text().strip()
                title = title_el.get_text().strip() if title_el else "No Title"
                link = title_el['href'] if title_el and 'href' in title_el.attrs else ""
                if link.startswith("//"):
                    link = "https:" + link
                if "uddg=" in link:
                    parsed = urllib.parse.urlparse(link)
                    qs = urllib.parse.parse_qs(parsed.query)
                    if "uddg" in qs:
                        link = qs["uddg"][0]
                results.append({"title": title, "snippet": snippet, "link": link, "source": "DuckDuckGo"})
            return results if results else None
    except Exception as e:
        print(f"DuckDuckGo HTML scrape error: {e}")
    return None

def search_duckduckgo_lite(query):
    try:
        import urllib.parse
        from bs4 import BeautifulSoup
        url = "https://lite.duckduckgo.com/lite/"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        r = requests.post(url, data={"q": query}, headers=headers, timeout=8)
        if r.status_code == 200:
            soup = BeautifulSoup(r.text, 'html.parser')
            results = []
            snippets = soup.find_all('td', class_='result-snippet')
            for snippet_td in snippets[:3]:
                parent_tr = snippet_td.find_parent('tr')
                if parent_tr:
                    prev_tr = parent_tr.find_previous_sibling('tr')
                    if prev_tr:
                        title_link = prev_tr.find('a', class_='result-link')
                        if title_link:
                            title = title_link.get_text().strip()
                            link = title_link['href']
                            if "uddg=" in link:
                                parsed = urllib.parse.urlparse(link)
                                qs = urllib.parse.parse_qs(parsed.query)
                                if "uddg" in qs:
                                    link = qs["uddg"][0]
                            snippet = snippet_td.get_text().strip()
                            results.append({"title": title, "snippet": snippet, "link": link, "source": "DuckDuckGo (Lite)"})
            return results if results else None
    except Exception as e:
        print(f"DuckDuckGo Lite scrape error: {e}")
    return None

def search_wikipedia(query):
    try:
        wiki_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{requests.utils.quote(query)}"
        wiki_response = requests.get(wiki_url, timeout=8)
        if wiki_response.status_code == 200:
            wiki_data = wiki_response.json()
            return [{"title": wiki_data.get("title", query), "snippet": wiki_data.get("extract", "")[:500], "source": "Wikipedia", "link": wiki_data.get("content_urls", {}).get("desktop", {}).get("page", "")}]
        return None
    except: return None

def search_anysearch(query):
    try:
        url = "https://anysearch-mcp.khulnasoft.com/search"
        response = requests.post(url, json={"query": query, "limit": 3}, timeout=10)
        if response.status_code == 200:
            data = response.json()
            results = []
            for item in data.get("results", [])[:3]:
                results.append({"title": item.get("title", ""), "snippet": item.get("snippet", ""), "link": item.get("url", ""), "source": "AnySearch"})
            return results if results else None
        return None
    except: return None

def search_web(query):
    all_results = []
    
    user_lower = query.lower()
    words = user_lower.split()
    is_university_query = 'seu' in words or 'seusl' in words or 'south eastern university' in user_lower or 'seu.ac.lk' in user_lower or 'seu ac lk' in user_lower
    
    search_query = query
    if is_university_query:
        # Simplify the query to prevent search engine confusion
        simplified = query.lower()
        noises = [
            'what is the', 'what is', 'who is the', 'who is', 'email address of', 'email of', 
            'email address', 'contact details of', 'contact of', 'phone number of', 'phone of', 
            'address of', 'dean of faculty of', 'dean of', 'faculty of', 'professor', 'sir', 
            'srilanka', 'sri lanka', 'south eastern university', 'seusl', 'seu'
        ]
        for word in noises:
            simplified = simplified.replace(word, '')
        
        # Remove leftover small connector words like 'of', 'for', 'the', 'and', 'in', 'at'
        words_list = [w for w in simplified.split() if w not in ['of', 'for', 'the', 'and', 'in', 'at', 'is', 'a']]
        simplified = " ".join(words_list)
        
        # If the user asked for contact details, ensure the keywords are present
        has_contact_req = any(w in query.lower() for w in ['email', 'mail', 'contact', 'phone', 'tele', 'number', 'address'])
        if has_contact_req:
            search_query = f"site:seu.ac.lk {simplified} email contact"
        else:
            search_query = f"site:seu.ac.lk {simplified}"
        
    ddg_results = search_duckduckgo(search_query)
    if ddg_results:
        all_results.extend(ddg_results)
    else:
        # Fallback to DDG Lite search
        lite_results = search_duckduckgo_lite(search_query)
        if lite_results:
            all_results.extend(lite_results)
            
    if not all_results and is_university_query:
        # Fallback to keyword-based search if strict domain search yields nothing
        fallback_results = search_duckduckgo(f"seu.ac.lk {query}")
        if fallback_results:
            all_results.extend(fallback_results)
        else:
            lite_fallback = search_duckduckgo_lite(f"seu.ac.lk {query}")
            if lite_fallback:
                all_results.extend(lite_fallback)
            
    # Wikipedia is only useful for general terms, not university-specific pages
    if not is_university_query:
        wiki_results = search_wikipedia(query)
        if wiki_results:
            all_results.extend(wiki_results)
            
    if len(all_results) < 2:
        any_results = search_anysearch(search_query)
        if any_results:
            all_results.extend(any_results)
        elif is_university_query:
            fallback_any = search_anysearch(f"seu.ac.lk {query}")
            if fallback_any:
                all_results.extend(fallback_any)
    
    # Scrape the top 2 results to get full page context
    if all_results:
        from bs4 import BeautifulSoup
        for r in all_results[:2]:
            link = r.get('link')
            if link and any(dom in link for dom in ['seu.ac.lk', 'wikipedia.org']):
                try:
                    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
                    page_response = requests.get(link, headers=headers, timeout=4)
                    if page_response.status_code == 200:
                        page_soup = BeautifulSoup(page_response.text, 'html.parser')
                        for element in page_soup(["script", "style", "nav", "footer", "header"]):
                            element.decompose()
                        page_text = page_soup.get_text()
                        clean_lines = [l.strip() for l in page_text.splitlines() if len(l.strip()) > 10]
                        page_content = "\n".join(clean_lines[:40])
                        if page_content:
                            r['scraped_content'] = page_content
                except Exception as e:
                    print(f"Error scraping search result link {link}: {e}")
                
    return all_results if all_results else None

def analyze_sentiment(text):
    positive_words = ['thanks', 'great', 'awesome', 'good', 'love', 'excellent', 'wonderful', 'perfect', 'helpful', 'amazing']
    negative_words = ['bad', 'terrible', 'awful', 'hate', 'useless', 'stupid', 'wrong', 'poor', 'frustrating', 'confusing']
    frustrated_words = ['confused', 'dont understand', 'not clear', 'what do you mean', 'explain again', 'still dont get']
    text_lower = text.lower()
    pos_count = sum(1 for w in positive_words if w in text_lower)
    neg_count = sum(1 for w in negative_words if w in text_lower)
    frus_count = sum(1 for w in frustrated_words if w in text_lower)
    if frus_count > 0: return 'frustrated'
    if neg_count > pos_count: return 'negative'
    if pos_count > 0: return 'positive'
    return 'neutral'

def needs_real_time_info(user_message):
    user_lower = user_message.lower()
    
    # Force real-time search if the query is about South Eastern University of Sri Lanka (SEUSL)
    words = user_lower.split()
    if 'seu' in words or 'seusl' in words or 'south eastern university' in user_lower or 'seu.ac.lk' in user_lower or 'seu ac lk' in user_lower:
        return True
        
    real_time_indicators = ['current', 'latest', 'today', 'now', '2024', '2025', '2026', 'president', 'prime minister', 'election', 'news', 'recent', 'weather', 'stock', 'price', 'score', 'live', 'update', 'who is the', 'who is current', 'currently', 'right now', 'who is president', 'who is prime minister', 'who leads', 'leader of', 'head of state', 'who governs', 'what is the capital', 'population of', 'weather in']
    return any(indicator in user_lower for indicator in real_time_indicators)

def generate_suggestions(user_message, response_text):
    user_lower = user_message.lower()
    if 'what is' in user_lower or 'define' in user_lower:
        return ["Can you give an example?", "Why is this important?", "How does this relate to my course?"]
    elif 'how' in user_lower:
        return ["Can you explain step by step?", "What are the prerequisites?", "Are there any alternatives?"]
    else:
        return ["Can you explain more?", "What's an example of this?", "How is this applied in practice?"]

# ========== PAST PAPER INTELLIGENCE ==========

past_paper_analysis = {}

def extract_questions_from_text(text):
    questions = []
    question_patterns = [
        r'(?:^|\n)\s*(?:\d+[\.\)]\s*)(.*?\?)',
        r'(?:^|\n)\s*(?:Q\d+[\.\):]\s*)(.*?\?)',
        r'(?:^|\n)\s*(?:Question\s*\d+[\.\):]\s*)(.*?\?)',
        r'(?:^|\n)\s*(?:[A-Z][\.\)]\s*)(.*?\?)',
        r'(?:^|\n)\s*(?:•|\-|\*)\s*(.*?\?)',
    ]
    for pattern in question_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
        questions.extend([q.strip() for q in matches if len(q.strip()) > 10])
    essay_patterns = [
        r'(?:^|\n)\s*(?:\d+[\.\)]\s*)((?:Discuss|Explain|Describe|Analyze|Compare|Evaluate|Define|Outline|Summarize|Justify|Illustrate).*?)(?:\n|$)',
        r'(?:^|\n)\s*(?:Q\d+[\.\):]\s*)((?:Discuss|Explain|Describe|Analyze|Compare|Evaluate|Define|Outline|Summarize|Justify|Illustrate).*?)(?:\n|$)',
    ]
    for pattern in essay_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
        questions.extend([q.strip() for q in matches if len(q.strip()) > 15 and q.strip() not in questions])
    seen = set()
    unique_questions = []
    for q in questions:
        if q.lower() not in seen:
            seen.add(q.lower())
            unique_questions.append(q)
    return unique_questions[:50]

def identify_topics_from_questions(questions):
    topic_keywords = {
        'Project Management': ['project', 'stakeholder', 'milestone', 'deliverable', 'gantt', 'wbs', 'scope', 'risk management'],
        'Software Development': ['software', 'development', 'agile', 'scrum', 'waterfall', 'sprint', 'testing', 'debugging'],
        'Database': ['database', 'sql', 'normalization', 'erd', 'query', 'table', 'index', 'transaction'],
        'Programming': ['programming', 'code', 'algorithm', 'function', 'variable', 'class', 'object', 'inheritance'],
        'Networking': ['network', 'protocol', 'tcp', 'ip', 'router', 'switch', 'firewall', 'dns'],
        'Security': ['security', 'encryption', 'authentication', 'firewall', 'vulnerability', 'threat', 'attack'],
        'Data Structures': ['data structure', 'array', 'linked list', 'tree', 'graph', 'stack', 'queue', 'hash'],
        'Operating Systems': ['operating system', 'process', 'thread', 'memory', 'scheduling', 'deadlock'],
        'Web Development': ['web', 'html', 'css', 'javascript', 'server', 'client', 'api', 'rest'],
        'AI & Machine Learning': ['artificial intelligence', 'machine learning', 'neural', 'algorithm', 'training', 'model'],
        'Requirements Engineering': ['requirement', 'specification', 'use case', 'user story', 'functional', 'non-functional'],
        'System Design': ['system design', 'architecture', 'scalability', 'performance', 'component', 'module'],
        'Testing': ['testing', 'unit test', 'integration', 'verification', 'validation', 'quality'],
        'Object-Oriented': ['object-oriented', 'oop', 'encapsulation', 'polymorphism', 'abstraction'],
    }
    topic_counts = {}
    for question in questions:
        question_lower = question.lower()
        for topic, keywords in topic_keywords.items():
            matching_keywords = [kw for kw in keywords if kw in question_lower]
            if matching_keywords:
                if topic not in topic_counts:
                    topic_counts[topic] = {'count': 0, 'keywords': set()}
                topic_counts[topic]['count'] += 1
                topic_counts[topic]['keywords'].update(matching_keywords)
    sorted_topics = sorted(topic_counts.items(), key=lambda x: x[1]['count'], reverse=True)
    return [
        {
            'topic': topic,
            'frequency': data['count'],
            'keywords': list(data['keywords']),
            'importance': 'High' if data['count'] >= 3 else ('Medium' if data['count'] >= 2 else 'Low')
        }
        for topic, data in sorted_topics
    ]

def analyze_past_paper(file_path, filename):
    file_type = filename.split('.')[-1].lower()
    if file_type == 'pdf':
        text = extract_text_from_pdf(file_path)
    elif file_type in ['pptx', 'ppt']:
        text = extract_text_from_pptx(file_path)
    elif file_type == 'docx':
        text = extract_text_from_docx(file_path)
    elif file_type == 'txt':
        text = extract_text_from_txt(file_path)
    else:
        text = extract_text_from_txt(file_path)
    text = clean_text(text)
    if not text or len(text.strip()) < 100:
        return None
    questions = extract_questions_from_text(text)
    topics = identify_topics_from_questions(questions)
    year_match = re.search(r'(20\d{2})', filename)
    year = year_match.group(1) if year_match else "Unknown"
    subject = filename.replace('.pdf', '').replace('.docx', '').replace('.txt', '')
    subject = re.sub(r'\d{4}', '', subject).strip('_ -')
    return {
        'filename': filename,
        'year': year,
        'subject': subject,
        'total_questions': len(questions),
        'questions': questions[:10],
        'topics': topics,
        'analyzed_at': datetime.utcnow().isoformat()
    }

def get_aggregated_topic_rankings():
    global past_paper_analysis
    all_topics = {}
    total_papers = len(past_paper_analysis)
    for filename, analysis in past_paper_analysis.items():
        for topic_data in analysis.get('topics', []):
            topic_name = topic_data['topic']
            if topic_name not in all_topics:
                all_topics[topic_name] = {
                    'topic': topic_name,
                    'total_frequency': 0,
                    'appearances': 0,
                    'keywords': set(),
                    'years': set()
                }
            all_topics[topic_name]['total_frequency'] += topic_data['frequency']
            all_topics[topic_name]['appearances'] += 1
            all_topics[topic_name]['keywords'].update(topic_data['keywords'])
            all_topics[topic_name]['years'].add(analysis.get('year', 'Unknown'))
    rankings = []
    for topic, data in all_topics.items():
        data['keywords'] = list(data['keywords'])
        data['years'] = sorted(list(data['years']))
        data['importance_score'] = (data['total_frequency'] * 0.6) + (data['appearances'] * 0.4)
        data['importance'] = 'High' if data['importance_score'] >= 5 else ('Medium' if data['importance_score'] >= 3 else 'Low')
        rankings.append(data)
    rankings.sort(key=lambda x: x['importance_score'], reverse=True)
    return {
        'total_papers_analyzed': total_papers,
        'total_questions_found': sum(len(analysis.get('questions', [])) for analysis in past_paper_analysis.values()),
        'topic_rankings': rankings
    }

# ========== SHARED PROMPT-BUILDING LOGIC ==========

def build_chat_context(user_message, session_id, length_control='medium', uploaded_context="", uploaded_filenames=None):
    user_lower = user_message.lower().strip()
    sentiment = analyze_sentiment(user_message)
    
    words = user_lower.split()
    is_university_query = 'seu' in words or 'seusl' in words or 'south eastern university' in user_lower or 'seu.ac.lk' in user_lower or 'seu ac lk' in user_lower

    follow_up_phrases = ['explain more', 'tell me more', 'give me more', 'elaborate', 'what about', 'can you explain', 'go deeper', 'more details', 'more explanation', 'expand', 'further', 'in detail', 'what else', 'continue', 'and then', 'why is that', 'how does that', 'can you clarify', 'what does that mean', 'explain it', 'describe it', 'tell about it', 'what is it', 'tell me about it', 'elaborate on that', 'go on']
    is_follow_up = any(phrase in user_lower for phrase in follow_up_phrases) and len(user_message.split()) <= 8

    history = get_session_history(session_id)
    recent_history = ""
    previous_topic = ""
    history_messages = []

    if history:
        last_messages = history[-6:]
        for msg in last_messages:
            role = "User" if msg['role'] == 'user' else "Assistant"
            recent_history += f"{role}: {msg['content']}\n"
        last_user_msgs = [m['content'] for m in history if m['role'] == 'user']
        if last_user_msgs: previous_topic = last_user_msgs[-1]
        history_messages = history[-10:]

    greetings = ['hi', 'hello', 'hey', 'good morning', 'good afternoon', 'good evening', 'sup', 'yo', 'hola', 'hii', 'heyy', 'helloo', 'morning', 'evening', 'good day']
    is_greeting = any(user_lower == g or user_lower.startswith(g + ' ') for g in greetings) and len(user_message.split()) <= 3

    identity_q = ['who are you', 'what are you', 'your name', 'about yourself', 'introduce yourself', 'tell me about yourself', 'who created you', 'are you ai', 'are you human', 'are you real']
    location_q = ['where are you', 'where do you live', 'your country', 'which country', 'where you from', 'your location']
    user_personal_q = ['know my name', 'do you know me', 'who am i', 'what is my name', 'remember me']
    is_identity = any(q in user_lower for q in identity_q)
    is_location = any(q in user_lower for q in location_q)
    is_user_personal = any(q in user_lower for q in user_personal_q)
    is_about_ai = is_identity or is_location

    thanks_words = ['thank', 'thanks', 'thx', 'appreciate']
    is_thanks = any(t in user_lower for t in thanks_words) and len(user_message.split()) <= 4

    needs_realtime = needs_real_time_info(user_message)
    is_casual = is_greeting or is_thanks or is_about_ai or is_user_personal

    web_results = None
    sources = []
    web_context = ""
    if needs_realtime and not is_casual:
        web_results = search_web(user_message)
        if web_results:
            sources = [r.get('title', '') for r in web_results[:3] if r.get('title')]
            web_context_items = []
            for r in web_results[:3]:
                title = r.get('title', '')
                snippet = r.get('snippet', '')
                scraped = r.get('scraped_content', '')
                item = f"📰 {title}:\nSnippet: {snippet}"
                if scraped:
                    item += f"\nFull Content:\n{scraped}"
                web_context_items.append(item)
            web_context = "\n\n".join(web_context_items)

    doc_context = ""
    if qdrant_client and embedding_model and not is_casual:
        search_query = previous_topic if (is_follow_up and previous_topic) else user_message
        try:
            query_embedding = embedding_model.encode(search_query).tolist()
            search_results = qdrant_client.search(collection_name="university_notes", query_vector=query_embedding, limit=15)
            if search_results:
                formatted_texts = []
                doc_names = []
                for hit in search_results:
                    if hit.score < 0.38:
                        continue
                    text = hit.payload.get('text', '')
                    filename = hit.payload.get('filename', 'Unknown source')
                    formatted_texts.append(f"[{filename}]: {text}")
                    if hit.payload.get('filename'):
                        doc_names.append(hit.payload.get('filename'))
                if formatted_texts:
                    doc_context = "\n\n".join(formatted_texts)
                sources = list(dict.fromkeys(sources + doc_names))
        except Exception as e: print(f"Search error: {e}")

    # Add past paper context
    if past_paper_analysis and not is_casual:
        past_paper_context = ""
        query_lower = user_message.lower()
        for filename, analysis in past_paper_analysis.items():
            for topic_data in analysis.get('topics', []):
                if any(kw in query_lower for kw in topic_data.get('keywords', [])):
                    past_paper_context += f"\n📝 Past Paper Topic: {topic_data['topic']} (Importance: {topic_data['importance']}, Year: {analysis.get('year', 'Unknown')})\n"
                    for q in analysis.get('questions', [])[:2]:
                        if any(kw in q.lower() for kw in topic_data.get('keywords', [])):
                            past_paper_context += f"   Sample Q: {q}\n"
        if past_paper_context:
            doc_context = past_paper_context + "\n" + doc_context

    if uploaded_context:
        doc_context = uploaded_context + "\n" + doc_context
    if uploaded_filenames:
        sources = list(dict.fromkeys(sources + uploaded_filenames))

    if is_greeting:
        system_prompt = "You are a friendly AI. Give a SHORT greeting. 1 sentence only."
        user_prompt = f"User: {user_message}\n\nShort greeting:"; max_tokens = 30
    elif is_identity:
        system_prompt = "You are an AI assistant. In 1-2 short sentences, explain you're an AI created to help people learn."
        user_prompt = f"User: {user_message}\n\nShort response:"; max_tokens = 60
    elif is_location:
        system_prompt = "You are an AI. In 1 short sentence, explain you don't have a physical location."
        user_prompt = f"User: {user_message}\n\nShort response:"; max_tokens = 40
    elif is_user_personal:
        system_prompt = "In 1-2 short sentences, honestly say you don't know their name but you're happy to help."
        user_prompt = f"User: {user_message}\n\nShort response:"; max_tokens = 40
    elif is_thanks:
        system_prompt = "Respond to thanks in 1 very short, warm sentence."
        user_prompt = f"User: {user_message}\n\nShort response:"; max_tokens = 20
    elif sentiment == 'frustrated':
        system_prompt = "The user seems frustrated. Be extra patient and helpful. 2-3 sentences."
        user_prompt = f"User seems confused: {user_message}\n\nPatient, helpful response:"; max_tokens = 150
    elif is_university_query:
        system_prompt = "You are the AI Learning Assistant of South Eastern University of Sri Lanka (SEUSL). Answer the question using the provided official university database and web pages. Be helpful, accurate, and professional. Prioritize database references and document contents. Cite file sources when available."
        
        context_parts = []
        if doc_context:
            context_parts.append(f"University Database Reference:\n{doc_context}")
        if web_context:
            context_parts.append(f"Web Search Results:\n{web_context}")
            
        combined_context = "\n\n".join(context_parts)
        user_prompt = f"Official Context:\n{combined_context}\n\nQuestion: {user_message}\n\nAnswer:"
        max_tokens = 350
    elif is_follow_up and doc_context:
        context_parts = []
        context_parts.append(f"Reference Documents:\n{doc_context}")
        if web_context:
            context_parts.append(f"Web Search Results:\n{web_context}")
        combined_context = "\n\n".join(context_parts)
        system_prompt = f"You are a helpful AI tutor. The user is asking a FOLLOW-UP about: '{previous_topic}'. Expand on it. 3-5 sentences. Use the provided Reference Documents as your primary source, citing the source filename (e.g. [filename.pdf]) when available. If the answer is not in the documents, state that clearly."
        user_prompt = f"Context:\n{combined_context}\n\nUser: {user_message}\n\nExpand on '{previous_topic}':"; max_tokens = 200
    elif is_follow_up and not doc_context:
        system_prompt = f"You are a helpful AI tutor. The user is asking a FOLLOW-UP about: '{previous_topic}'. Expand on it. 3-5 sentences."
        user_prompt = f"Previous topic: '{previous_topic}'. User: {user_message}\n\nExpand on '{previous_topic}':"; max_tokens = 200
    elif doc_context:
        context_parts = []
        context_parts.append(f"Reference Documents:\n{doc_context}")
        if web_context:
            context_parts.append(f"Web Search Results:\n{web_context}")
        combined_context = "\n\n".join(context_parts)
        system_prompt = (
            "You are a helpful AI tutor. Answer in 1-3 SHORT sentences. Use the provided Reference Documents to answer the question, citing the source filename (e.g. [filename.pdf]) when available. "
            "Strictly prioritize the Reference Documents. If they do not contain the answer, state: 'Based on the uploaded documents, I cannot find the answer to this question.' Do not use general knowledge or web search if the query is about the uploaded documents."
        )
        user_prompt = f"Context:\n{combined_context}\n\nQuestion: {user_message}\n\nShort answer:"; max_tokens = 100
    elif web_results:
        system_prompt = "You are a helpful AI with web access. Answer in 1-3 SHORT sentences. Be direct."
        user_prompt = f"Web results:\n{web_context}\n\nQuestion: {user_message}\n\nShort answer:"; max_tokens = 120
    else:
        system_prompt = "You are a smart AI assistant. Answer in 1-3 SHORT sentences."
        user_prompt = f"Question: {user_message}\n\nShort answer:"; max_tokens = 100

    if length_control == 'short':
        system_prompt = system_prompt.replace("3-5 sentences", "1 sentence max").replace("2-3 sentences", "1 sentence max")
        system_prompt += " IMPORTANT: Keep response extremely short and concise (1 sentence max)."
        max_tokens = 80
    elif length_control == 'detailed':
        system_prompt = system_prompt.replace("Answer in 1-3 SHORT sentences.", "Provide a detailed and thorough explanation.")
        system_prompt = system_prompt.replace("SHORT greeting. 1 sentence only.", "friendly greeting.")
        system_prompt = system_prompt.replace("In 1-2 short sentences, explain", "Explain in detail")
        system_prompt = system_prompt.replace("In 1 short sentence, explain", "Explain in detail")
        system_prompt = system_prompt.replace("1-2 short sentences, honestly say", "Honestly say")
        system_prompt = system_prompt.replace("1 very short, warm sentence.", "warm response.")
        system_prompt = system_prompt.replace("2-3 sentences.", "detailed explanation.")
        system_prompt = system_prompt.replace("3-5 sentences.", "comprehensive, detailed explanation.")
        
        user_prompt = user_prompt.replace("Short answer:", "Detailed answer:")
        user_prompt = user_prompt.replace("Short response:", "Detailed response:")
        user_prompt = user_prompt.replace("Short greeting:", "Greeting:")
        
        system_prompt += " Provide a detailed explanation with formatting, bullet points, or code blocks where relevant."
        max_tokens = 1000
    else:
        system_prompt = system_prompt.replace("Answer in 1-3 SHORT sentences.", "Provide a clear and complete answer.")
        system_prompt = system_prompt.replace("SHORT greeting. 1 sentence only.", "warm greeting.")
        system_prompt = system_prompt.replace("In 1-2 short sentences, explain", "Explain")
        system_prompt = system_prompt.replace("In 1 short sentence, explain", "Explain")
        system_prompt = system_prompt.replace("1-2 short sentences, honestly say", "Honestly say")
        system_prompt = system_prompt.replace("1 very short, warm sentence.", "warm response.")
        system_prompt = system_prompt.replace("2-3 sentences.", "complete response.")
        system_prompt = system_prompt.replace("3-5 sentences.", "clear response.")
        
        user_prompt = user_prompt.replace("Short answer:", "Answer:")
        user_prompt = user_prompt.replace("Short response:", "Response:")
        user_prompt = user_prompt.replace("Short greeting:", "Greeting:")
        
        system_prompt += " Keep response to a medium length. Use bullet points or lists if the source documents contain structured information or lists."
        max_tokens = 500

    wants_chart = any(w in user_lower for w in ['chart', 'graph', 'plot', 'statistics', 'stats', 'compare', 'comparison', 'data table', 'trend'])
    if wants_chart and not is_casual:
        system_prompt += (
            " If the user asks for comparison, statistics, or trends, format the numeric data as an interactive chart at the end of your message in this exact format: "
            "[CHART: {\"type\": \"bar\", \"title\": \"Chart Title\", \"labels\": [\"Label1\", \"Label2\"], \"data\": [10, 20]}]. "
            "Support chart types: 'bar', 'line', 'pie', 'doughnut'. Keep labels short and data values simple integers/numbers."
        )
        max_tokens = max(max_tokens, 600)

    return {
        'system_prompt': system_prompt, 'user_prompt': user_prompt, 'max_tokens': max_tokens,
        'mode': 'followup' if is_follow_up else ('realtime' if web_results else ('study' if doc_context else 'general')),
        'sentiment': sentiment, 'is_greeting': is_greeting, 'is_thanks': is_thanks,
        'is_follow_up': is_follow_up, 'history_messages': history_messages, 'sources': sources,
    }

def postprocess_response(text, ctx):
    if not text: return "How can I help you today?"
    text = text.strip()
    if ctx['is_greeting'] and (len(text) < 2 or len(text) > 60):
        text = "Hey there! 👋 How can I help you today?"
    if ctx['is_thanks'] and len(text) > 30:
        text = "You're welcome! 😊"
    return text if len(text) >= 2 else "How can I help you today?"

# ==================== ROUTES ====================

@app.route('/')
@login_required
def index():
    return render_template('index.html')

# ========== STUDENT AUTH (MongoDB) ==========

@app.route('/login', methods=['GET', 'POST'])
def login_page():
    if request.method == 'POST':
        data = request.json or {}
        email = data.get('email', '').strip().lower()
        password = data.get('password', '')
        if not email or not password:
            return jsonify({'success': False, 'message': 'Email and password required'}), 400
        
        if users_collection is None:
            # SQLite fallback
            try:
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute('SELECT id, name, email, password_hash FROM users WHERE email = ?', (email,))
                user_row = cursor.fetchone()
                conn.close()
                if user_row and check_password_hash(user_row[3], password):
                    session.permanent = True
                    session['student_user'] = {'id': str(user_row[0]), 'email': user_row[2], 'name': user_row[1]}
                    return jsonify({'success': True})
            except Exception as e:
                print(f"SQLite login error: {e}")
                return jsonify({'success': False, 'message': 'Database error'}), 500
            return jsonify({'success': False, 'message': 'Invalid email or password'}), 401
            
        try:
            user = users_collection.find_one({'email': email})
            if user and check_password_hash(user['password_hash'], password):
                session.permanent = True
                session['student_user'] = {'id': str(user['_id']), 'email': email, 'name': user.get('name', email)}
                return jsonify({'success': True})
        except Exception as e:
            print(f"MongoDB login error: {e}")
        return jsonify({'success': False, 'message': 'Invalid email or password'}), 401
    if 'student_user' in session:
        return redirect(url_for('index'))
    return render_template('login.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup_page():
    if request.method == 'POST':
        data = request.json or {}
        name = data.get('name', '').strip()
        email = data.get('email', '').strip().lower()
        password = data.get('password', '')
        if not name or not email or not password:
            return jsonify({'success': False, 'message': 'All fields required'}), 400
        if len(password) < 6:
            return jsonify({'success': False, 'message': 'Password must be 6+ characters'}), 400
            
        if users_collection is None:
            # SQLite fallback
            try:
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute('SELECT id FROM users WHERE email = ?', (email,))
                if cursor.fetchone():
                    conn.close()
                    return jsonify({'success': False, 'message': 'Email already registered'}), 400
                
                cursor.execute(
                    'INSERT INTO users (name, email, password_hash, created_at) VALUES (?, ?, ?, ?)',
                    (name, email, generate_password_hash(password), datetime.utcnow().isoformat())
                )
                conn.commit()
                conn.close()
                print(f"✅ New student (SQLite): {name} ({email})")
                return jsonify({'success': True, 'message': 'Registration successful! Please login.'})
            except Exception as e:
                print(f"SQLite signup error: {e}")
                return jsonify({'success': False, 'message': 'Registration failed.'}), 500
                
        try:
            if users_collection.find_one({'email': email}):
                return jsonify({'success': False, 'message': 'Email already registered'}), 400
            users_collection.insert_one({'name': name, 'email': email, 'password_hash': generate_password_hash(password), 'created_at': datetime.utcnow()})
            print(f"✅ New student: {name} ({email})")
            return jsonify({'success': True, 'message': 'Registration successful! Please login.'})
        except Exception as e:
            print(f"MongoDB signup error: {e}")
            return jsonify({'success': False, 'message': 'Registration failed.'}), 500
    if 'student_user' in session:
        return redirect(url_for('index'))
    return render_template('signup.html')

@app.route('/logout')
def logout():
    session.pop('student_user', None)
    return redirect(url_for('login_page'))

# ========== CHAT SESSION MANAGEMENT ==========

@app.route('/api/sessions', methods=['GET'])
@login_required
def api_get_sessions():
    if not mongo_available or sessions_collection is None:
        return jsonify({'success': True, 'sessions': []})
    try:
        docs = sessions_collection.find({'user_email': session['student_user']['email']})
        sessions_data = []
        for doc in docs:
            messages = doc.get('messages', [])
            sessions_data.append({'id': doc['session_id'], 'title': doc.get('title', 'New Chat'), 'messages': [{'role': m['role'], 'content': m['content']} for m in messages], 'timestamp': str(doc.get('_id').generation_time) if doc.get('_id') else ''})
        return jsonify({'success': True, 'sessions': sessions_data})
    except Exception as e:
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/sessions', methods=['POST'])
@login_required
def api_create_session():
    data = request.json or {}
    session_id = data.get('id')
    title = data.get('title', 'New Chat')
    if not session_id:
        return jsonify({'success': False, 'message': 'Session ID required'}), 400
    if mongo_available and sessions_collection is not None:
        try:
            sessions_collection.update_one({'session_id': session_id}, {'$set': {'user_email': session['student_user']['email'], 'title': title}}, upsert=True)
        except: pass
    return jsonify({'success': True})

@app.route('/api/sessions/update-title', methods=['POST'])
@login_required
def api_update_session_title():
    data = request.json or {}
    session_id = data.get('id')
    title = data.get('title')
    if not session_id or not title:
        return jsonify({'success': False, 'message': 'Session ID and title required'}), 400
    if mongo_available and sessions_collection is not None:
        try:
            sessions_collection.update_one({'session_id': session_id}, {'$set': {'title': title}})
        except: pass
    return jsonify({'success': True})

@app.route('/api/sessions/<session_id>', methods=['DELETE'])
@login_required
def api_delete_session(session_id):
    if mongo_available and sessions_collection is not None:
        try:
            sessions_collection.delete_one({'session_id': session_id})
        except: pass
    return jsonify({'success': True})

# ========== ADMIN ROUTES ==========

@app.route('/admin')
def admin_login_page():
    if 'admin_logged_in' in session:
        return redirect(url_for('admin_panel'))
    return render_template('admin_login.html')

@app.route('/admin/login', methods=['POST'])
def admin_login():
    data = request.json
    if data.get('username') == ADMIN_USERNAME and check_password_hash(ADMIN_PASSWORD_HASH, data.get('password', '')):
        session.permanent = True
        session['admin_logged_in'] = True
        return jsonify({'success': True})
    return jsonify({'success': False, 'message': 'Invalid credentials'}), 401

@app.route('/admin/logout')
def admin_logout():
    session.pop('admin_logged_in', None)
    return redirect(url_for('admin_login_page'))

@app.route('/admin/panel')
@admin_required
def admin_panel():
    return render_template('admin.html')

@app.route('/admin/upload', methods=['POST'])
@admin_required
def upload_file():
    if 'files' not in request.files:
        return jsonify({'error': 'No files'}), 400
    files = request.files.getlist('files')
    category = request.form.get('category', '')
    uploaded, failed = [], []
    for file in files:
        if not file.filename: continue
        try:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)
            file_type = filename.split('.')[-1].lower()
            if file_type == 'pdf': text = extract_text_from_pdf(file_path)
            elif file_type in ['pptx', 'ppt']: text = extract_text_from_pptx(file_path)
            elif file_type == 'docx': text = extract_text_from_docx(file_path)
            elif file_type == 'txt': text = extract_text_from_txt(file_path)
            elif file_type in ['csv', 'xlsx', 'xls']:
                try: text = pd.read_csv(file_path).to_string() if file_type == 'csv' else pd.read_excel(file_path).to_string()
                except: text = ""
            else: text = extract_text_from_txt(file_path)
            text = clean_text(text)
            if text and len(text.strip()) > 50:
                chunks = chunk_text(text)
                points = []
                for i, chunk in enumerate(chunks):
                    embedding = embedding_model.encode(chunk).tolist()
                    points.append(PointStruct(id=str(uuid.uuid4()), vector=embedding, payload={"filename": filename, "text": chunk, "chunk_index": i, "category": category, "file_type": file_type, "upload_date": str(pd.Timestamp.now())}))
                if qdrant_client: qdrant_client.upsert(collection_name="university_notes", points=points)
                uploaded.append({'name': filename, 'type': file_type.upper(), 'chunks': len(chunks)})
            else: failed.append({'name': filename, 'reason': 'No text extracted'})
            if os.path.exists(file_path): os.remove(file_path)
        except Exception as e: failed.append({'name': file.filename, 'reason': str(e)})
    return jsonify({'success': True, 'uploaded': uploaded, 'failed': failed})

# ========== PAST PAPER UPLOAD ==========

@app.route('/admin/upload-past-paper', methods=['POST'])
@admin_required
def upload_past_paper():
    global past_paper_analysis
    if 'files' not in request.files:
        return jsonify({'error': 'No files'}), 400
    files = request.files.getlist('files')
    results = {'uploaded': [], 'failed': []}
    for file in files:
        if not file.filename:
            continue
        try:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"past_paper_{uuid.uuid4()}_{filename}")
            file.save(file_path)
            analysis = analyze_past_paper(file_path, filename)
            if analysis:
                past_paper_analysis[filename] = analysis
                results['uploaded'].append({
                    'name': filename,
                    'year': analysis['year'],
                    'questions_found': analysis['total_questions'],
                    'topics_found': len(analysis['topics'])
                })
            else:
                results['failed'].append({'name': filename, 'reason': 'Could not extract sufficient text'})
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            results['failed'].append({'name': file.filename, 'reason': str(e)})
    return jsonify({
        'success': True,
        'results': results,
        'total_papers_analyzed': len(past_paper_analysis)
    })

# ========== PAST PAPER API ==========

@app.route('/api/past-papers/rankings', methods=['GET'])
@login_required
def get_topic_rankings():
    rankings = get_aggregated_topic_rankings()
    return jsonify({'success': True, 'rankings': rankings})

@app.route('/api/past-papers/list', methods=['GET'])
@login_required
def list_past_papers():
    papers = []
    for filename, analysis in past_paper_analysis.items():
        papers.append({
            'filename': filename,
            'year': analysis.get('year', 'Unknown'),
            'subject': analysis.get('subject', 'General'),
            'questions_count': analysis.get('total_questions', 0),
            'topics_count': len(analysis.get('topics', [])),
            'analyzed_at': analysis.get('analyzed_at', '')
        })
    return jsonify({'success': True, 'papers': sorted(papers, key=lambda x: x['year'], reverse=True)})

@app.route('/api/past-papers/search', methods=['POST'])
@login_required
def search_past_papers():
    data = request.json or {}
    query = data.get('query', '').strip().lower()
    if not query:
        return jsonify({'success': False, 'message': 'No query provided'}), 400
    relevant_questions = []
    relevant_topics = []
    for filename, analysis in past_paper_analysis.items():
        for question in analysis.get('questions', []):
            if query in question.lower():
                relevant_questions.append({
                    'question': question,
                    'paper': filename,
                    'year': analysis.get('year', 'Unknown')
                })
        for topic_data in analysis.get('topics', []):
            if query in topic_data['topic'].lower() or any(query in kw.lower() for kw in topic_data.get('keywords', [])):
                relevant_topics.append({
                    'topic': topic_data['topic'],
                    'frequency': topic_data['frequency'],
                    'importance': topic_data['importance'],
                    'paper': filename,
                    'year': analysis.get('year', 'Unknown')
                })
    return jsonify({
        'success': True,
        'query': query,
        'related_questions': relevant_questions[:10],
        'related_topics': relevant_topics[:5]
    })

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    try:
        uploaded_context, uploaded_filenames = "", []
        if request.content_type and 'multipart/form-data' in request.content_type:
            user_message = request.form.get('message', '').strip()
            session_id = request.form.get('session_id', 'default')
            length_control = request.form.get('length_control', 'medium')
            files = request.files.getlist('files')
        else:
            data = request.get_json(silent=True) or {}
            user_message = data.get('message', '').strip()
            session_id = data.get('session_id', 'default')
            length_control = data.get('length_control', 'medium')
            files = []
        for file in files:
            if not file.filename: continue
            try:
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_{uuid.uuid4()}_{filename}")
                file.save(file_path)
                file_type = filename.split('.')[-1].lower()
                if file_type == 'pdf': file_text = extract_text_from_pdf(file_path)
                elif file_type == 'pptx': file_text = extract_text_from_pptx(file_path)
                elif file_type == 'docx': file_text = extract_text_from_docx(file_path)
                elif file_type == 'txt': file_text = extract_text_from_txt(file_path)
                else: file_text = extract_text_from_txt(file_path)
                file_text = clean_text(file_text)
                if file_text:
                    uploaded_context += f"--- Content from uploaded file '{filename}' ---\n{file_text}\n\n"
                    uploaded_filenames.append(filename)
                if os.path.exists(file_path): os.remove(file_path)
            except Exception as e: print(f"Error processing temp upload: {e}")
        if not user_message and uploaded_filenames:
            user_message = f"Describe the attached file: {', '.join(uploaded_filenames)}"
        if not user_message: return jsonify({'response': 'Please type a message.'}), 400
        if not gemini_api_key and not groq_api_key: return jsonify({'response': 'AI service not available.'}), 500
        ctx = build_chat_context(user_message, session_id, length_control, uploaded_context, uploaded_filenames)
        response_text = None
        if gemini_api_key:
            contents = build_gemini_contents(ctx.get('history_messages', []), ctx['user_prompt'])
            for model in GEMINI_MODELS:
                try:
                    response_text = gemini_chat_completion(system_prompt=ctx['system_prompt'], contents=contents, model=model, max_tokens=ctx['max_tokens'], temperature=0.7)
                    if response_text: break
                except: continue
        if not response_text and groq_api_key:
            messages = [{"role": "system", "content": ctx['system_prompt']}]
            if ctx.get('history_messages'): messages.extend(ctx['history_messages'])
            messages.append({"role": "user", "content": ctx['user_prompt']})
            for model in GROQ_MODELS:
                try:
                    response_text = groq_chat_completion(messages=messages, model=model, max_tokens=ctx['max_tokens'], temperature=0.7)
                    break
                except: continue
        if response_text:
            response_text = postprocess_response(response_text, ctx)
            save_message(session_id, 'user', user_message)
            save_message(session_id, 'assistant', response_text)
            trim_session_history(session_id)
            suggestions = generate_suggestions(user_message, response_text)
            return jsonify({'response': response_text, 'sources': ctx.get('sources', []), 'mode': ctx['mode'], 'sentiment': ctx['sentiment'], 'suggestions': suggestions})
        else:
            return jsonify({'response': "I'm having trouble. Try again!"}), 500
    except Exception as e:
        print(f"Chat error: {e}")
        return jsonify({'response': "Sorry, an error occurred."}), 500

@app.route('/chat/stream', methods=['POST'])
@login_required
def chat_stream():
    uploaded_context, uploaded_filenames = "", []
    image_data, image_mime = None, None
    if request.content_type and 'multipart/form-data' in request.content_type:
        user_message = request.form.get('message', '').strip()
        session_id = request.form.get('session_id', 'default')
        length_control = request.form.get('length_control', 'medium')
        files = request.files.getlist('files')
    else:
        data = request.get_json(silent=True) or {}
        user_message = data.get('message', '').strip()
        session_id = data.get('session_id', 'default')
        length_control = data.get('length_control', 'medium')
        image_data = data.get('image_data')
        image_mime = data.get('image_mime', 'image/png')
        files = []
    for file in files:
        if not file.filename: continue
        try:
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_{uuid.uuid4()}_{filename}")
            file.save(file_path)
            file_type = filename.split('.')[-1].lower()
            if file_type == 'pdf': file_text = extract_text_from_pdf(file_path)
            elif file_type == 'pptx': file_text = extract_text_from_pptx(file_path)
            elif file_type == 'docx': file_text = extract_text_from_docx(file_path)
            elif file_type == 'txt': file_text = extract_text_from_txt(file_path)
            else: file_text = extract_text_from_txt(file_path)
            file_text = clean_text(file_text)
            if file_text:
                uploaded_context += f"--- Content from uploaded file '{filename}' ---\n{file_text}\n\n"
                uploaded_filenames.append(filename)
            if os.path.exists(file_path): os.remove(file_path)
        except Exception as e: print(f"Error processing temp upload: {e}")
    if not user_message and uploaded_filenames:
        user_message = f"Describe the attached file: {', '.join(uploaded_filenames)}"
    if not user_message: return jsonify({'error': 'Please type a message.'}), 400
    if not gemini_api_key and not groq_api_key: return jsonify({'error': 'AI service not available.'}), 500
    ctx = build_chat_context(user_message, session_id, length_control, uploaded_context, uploaded_filenames)
    def event_stream():
        full_response = ""
        streamed_ok = False
        if gemini_api_key:
            contents = build_gemini_contents(ctx.get('history_messages', []), ctx['user_prompt'])
            if image_data:
                # Append base64 image details to the last part (current user prompt)
                contents[-1]['parts'].append({
                    "inlineData": {
                        "mimeType": image_mime,
                        "data": image_data
                    }
                })
            for model in GEMINI_MODELS:
                try:
                    for delta in gemini_chat_completion_stream(system_prompt=ctx['system_prompt'], contents=contents, model=model, max_tokens=ctx['max_tokens'], temperature=0.7):
                        full_response += delta
                        yield f"data: {json.dumps({'delta': delta})}\n\n"
                    streamed_ok = True; break
                except: full_response = ""
        if not streamed_ok and groq_api_key:
            messages = [{"role": "system", "content": ctx['system_prompt']}]
            if ctx.get('history_messages'): messages.extend(ctx['history_messages'])
            
            if image_data:
                messages.append({
                    "role": "user",
                    "content": [
                        {"type": "text", "text": ctx['user_prompt']},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{image_mime};base64,{image_data}"
                            }
                        }
                    ]
                })
                models_to_try = ["llama-3.2-11b-vision-preview"] + GROQ_MODELS
            else:
                messages.append({"role": "user", "content": ctx['user_prompt']})
                models_to_try = GROQ_MODELS

            for model in models_to_try:
                try:
                    for delta in groq_chat_completion_stream(messages=messages, model=model, max_tokens=ctx['max_tokens'], temperature=0.7):
                        full_response += delta
                        yield f"data: {json.dumps({'delta': delta})}\n\n"
                    streamed_ok = True; break
                except: full_response = ""
        if not streamed_ok:
            yield f"data: {json.dumps({'error': 'AI service is having trouble.'})}\n\n"
            return
        clean_response = postprocess_response(full_response, ctx)
        save_message(session_id, 'user', user_message)
        save_message(session_id, 'assistant', clean_response)
        trim_session_history(session_id)
        suggestions = generate_suggestions(user_message, clean_response)
        yield f"data: {json.dumps({'done': True, 'mode': ctx['mode'], 'sentiment': ctx['sentiment'], 'suggestions': suggestions, 'sources': ctx.get('sources', [])})}\n\n"
    return Response(event_stream(), mimetype='text/event-stream', headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'})

# ========== SPECIAL FEATURES ==========

@app.route('/generate-flashcards', methods=['POST'])
@login_required
def generate_flashcards():
    try:
        data = request.json or {}
        custom_topic = data.get('topic', '').strip()
        if not gemini_api_key and not groq_api_key: return jsonify({'error': 'AI service not available.'}), 500
        doc_context = ""
        if qdrant_client and embedding_model:
            try:
                if custom_topic:
                    query_embedding = embedding_model.encode(custom_topic).tolist()
                    results = qdrant_client.search(collection_name="university_notes", query_vector=query_embedding, limit=4)
                    if results: doc_context = "\n\n".join([h.payload.get('text', '') for h in results if h.score >= 0.0])
                else:
                    results = qdrant_client.scroll(collection_name="university_notes", limit=5)
                    if results[0]: doc_context = "\n\n".join([h.payload.get('text', '') for h in results[0]])
            except Exception as e: print(f"Flashcards search error: {e}")
        system_prompt = "Generate exactly 10 flashcards. Respond ONLY with a valid JSON array: [{\"question\": \"...\", \"answer\": \"...\"}]"
        user_prompt = f"Generate 10 flashcards. Topic: {custom_topic or 'General'}"
        if doc_context: user_prompt += f"\n\nContext:\n{doc_context[:2000]}"
        response_text = None
        if gemini_api_key:
            for model in GEMINI_MODELS:
                try:
                    response_text = gemini_chat_completion(system_prompt=system_prompt, contents=[{"role":"user","parts":[{"text":user_prompt}]}], model=model, max_tokens=1000, temperature=0.7)
                    if response_text: break
                except: continue
        if not response_text and groq_api_key:
            for model in GROQ_MODELS:
                try:
                    response_text = groq_chat_completion(messages=[{"role":"system","content":system_prompt},{"role":"user","content":user_prompt}], model=model, max_tokens=1000, temperature=0.7)
                    break
                except: continue
        if not response_text: return jsonify({'error': 'Failed to generate flashcards.'}), 500
        raw = response_text.strip()
        if raw.startswith("```"): raw = re.sub(r'^```(?:json)?\n', '', raw); raw = re.sub(r'\n```$', '', raw)
        try:
            flashcards = json.loads(raw)
            if isinstance(flashcards, list) and len(flashcards) > 0:
                return jsonify({'success': True, 'flashcards': flashcards})
        except:
            try:
                match = re.search(r'\[\s*\{.*\}\s*\]', raw, re.DOTALL)
                if match: flashcards = json.loads(match.group(0)); return jsonify({'success': True, 'flashcards': flashcards})
            except: pass
        return jsonify({'success': True, 'flashcards': [{"question":"What is a project?","answer":"A temporary endeavor to create a unique product or service."}], 'note':'Using fallback cards'})
    except Exception as e:
        print(f"Flashcards error: {e}")
        return jsonify({'error': 'An internal error occurred.'}), 500

@app.route('/generate-summary', methods=['POST'])
@login_required
def generate_summary():
    try:
        data = request.json or {}
        topic = data.get('topic', '').strip()
        if not topic: return jsonify({'error': 'Please provide a topic.'}), 400
        if not gemini_api_key and not groq_api_key: return jsonify({'error': 'AI service not available.'}), 500
        doc_context = ""
        if qdrant_client and embedding_model:
            try:
                query_embedding = embedding_model.encode(topic).tolist()
                results = qdrant_client.search(collection_name="university_notes", query_vector=query_embedding, limit=6)
                if results: doc_context = "\n\n".join([h.payload.get('text', '') for h in results if h.score >= 0.0])
            except: pass
        if not doc_context: return jsonify({'error': 'No relevant documents found.'}), 404
        system_prompt = "You are an academic summarizer. Create a structured Markdown summary with headings, subheadings, and bullet points."
        user_prompt = f"Topic: {topic}\n\nNotes:\n{doc_context[:3000]}\n\nSummary:"
        response_text = None
        if gemini_api_key:
            for model in GEMINI_MODELS:
                try:
                    response_text = gemini_chat_completion(system_prompt=system_prompt, contents=[{"role":"user","parts":[{"text":user_prompt}]}], model=model, max_tokens=1500, temperature=0.5)
                    if response_text: break
                except: continue
        if not response_text and groq_api_key:
            for model in GROQ_MODELS:
                try:
                    response_text = groq_chat_completion(messages=[{"role":"system","content":system_prompt},{"role":"user","content":user_prompt}], model=model, max_tokens=1500, temperature=0.5)
                    break
                except: continue
        if response_text: return jsonify({'success': True, 'summary': response_text.strip()})
        return jsonify({'error': 'Failed to generate summary.'}), 500
    except Exception as e:
        print(f"Summary error: {e}")
        return jsonify({'error': 'An internal error occurred.'}), 500

@app.route('/generate-quiz', methods=['POST'])
@login_required
def generate_quiz():
    try:
        data = request.json or {}
        topic = data.get('topic', '').strip()
        if not topic: return jsonify({'error': 'Please provide a topic.'}), 400
        if not gemini_api_key and not groq_api_key: return jsonify({'error': 'AI service not available.'}), 500
        
        doc_context = ""
        if qdrant_client and embedding_model:
            try:
                query_embedding = embedding_model.encode(topic).tolist()
                results = qdrant_client.search(collection_name="university_notes", query_vector=query_embedding, limit=6)
                if results: doc_context = "\n\n".join([h.payload.get('text', '') for h in results if h.score >= 0.0])
            except: pass
            
        if not doc_context: return jsonify({'error': 'No relevant notes found for this topic.'}), 404
        
        system_prompt = (
            "You are an academic quiz generator. Create a structured multiple-choice quiz on the topic. "
            "Your output must be a VALID JSON array of 5 questions. "
            "Do not include markdown tags (like ```json) in your raw response. Output ONLY the raw JSON string. "
            "Format of each question object:\n"
            "[\n"
            "  {\n"
            "    \"question\": \"Question text?\",\n"
            "    \"options\": [\"Option A\", \"Option B\", \"Option C\", \"Option D\"],\n"
            "    \"answer_index\": 0,\n"
            "    \"explanation\": \"Explain why the option at answer_index is correct based on academic theory.\"\n"
            "  }\n"
            "]"
        )
        user_prompt = f"Topic: {topic}\n\nNotes Context:\n{doc_context[:3000]}\n\nGenerate 5 MCQs in raw JSON:"
        
        response_text = None
        if gemini_api_key:
            for model in GEMINI_MODELS:
                try:
                    response_text = gemini_chat_completion(system_prompt=system_prompt, contents=[{"role":"user","parts":[{"text":user_prompt}]}], model=model, max_tokens=1500, temperature=0.5)
                    if response_text: break
                except: continue
        if not response_text and groq_api_key:
            for model in GROQ_MODELS:
                try:
                    response_text = groq_chat_completion(messages=[{"role":"system","content":system_prompt},{"role":"user","content":user_prompt}], model=model, max_tokens=1500, temperature=0.5)
                    break
                except: continue
                
        if not response_text: return jsonify({'error': 'Failed to generate quiz.'}), 500
        
        raw = response_text.strip()
        if raw.startswith("```"): 
            raw = re.sub(r'^```(?:json)?\n', '', raw)
            raw = re.sub(r'\n```$', '', raw)
            
        try:
            quiz = json.loads(raw)
            if isinstance(quiz, list) and len(quiz) > 0:
                return jsonify({'success': True, 'quiz': quiz})
        except:
            try:
                match = re.search(r'\[\s*\{.*\}\s*\]', raw, re.DOTALL)
                if match: 
                    quiz = json.loads(match.group(0))
                    return jsonify({'success': True, 'quiz': quiz})
            except: pass
            
        fallback_quiz = [
            {
                "question": "What is the primary goal of project risk management?",
                "options": [
                    "To identify and manage risks to ensure project success",
                    "To completely eliminate all risks from the project",
                    "To transfer all risks to third parties",
                    "To increase the project budget"
                ],
                "answer_index": 0,
                "explanation": "Project risk management aims to identify, assess, and manage risks, minimizing negative impacts and maximizing positive opportunities for project success."
            }
        ]
        return jsonify({'success': True, 'quiz': fallback_quiz, 'note': 'Using fallback quiz'})
        
    except Exception as e:
        print(f"Quiz error: {e}")
        return jsonify({'error': 'An internal error occurred.'}), 500

# ========== SMART STUDY PLAN GENERATOR ==========

@app.route('/generate-study-plan', methods=['POST'])
@login_required
def generate_study_plan():
    try:
        data = request.json or {}
        exam_date = data.get('exam_date', '').strip()
        subjects_str = data.get('subjects', '').strip()
        study_hours_per_day = int(data.get('hours_per_day', 4))
        
        if not exam_date: return jsonify({'error': 'Please provide your exam date.'}), 400
        if not subjects_str: return jsonify({'error': 'Please list your subjects.'}), 400
        
        subjects = [s.strip() for s in subjects_str.split(',') if s.strip()]
        if not subjects: return jsonify({'error': 'Please enter at least one subject.'}), 400
        
        try:
            exam_datetime = datetime.strptime(exam_date, '%Y-%m-%d')
            today = datetime.utcnow()
            days_until_exam = (exam_datetime - today).days
            if days_until_exam <= 0: return jsonify({'error': 'Exam date must be in the future!'}), 400
            if days_until_exam > 365: return jsonify({'error': 'Exam date is too far away.'}), 400
        except ValueError:
            return jsonify({'error': 'Invalid date format. Use YYYY-MM-DD'}), 400
        
        total_study_hours = days_until_exam * study_hours_per_day
        subject_content = {}
        total_content_volume = 0
        
        if qdrant_client and embedding_model:
            for subject in subjects:
                try:
                    query_embedding = embedding_model.encode(subject).tolist()
                    results = qdrant_client.search(collection_name="university_notes", query_vector=query_embedding, limit=5)
                    if results:
                        texts = [h.payload.get('text', '') for h in results if h.score >= 0.0]
                        combined = "\n\n".join(texts)
                        word_count = len(combined.split())
                        subject_content[subject] = {'text': combined[:2000], 'word_count': word_count, 'chunks_found': len(texts)}
                        total_content_volume += word_count
                except: pass
        
        subjects_list = "\n".join([f"- {s}" for s in subjects])
        content_summary = ""
        for subj, info in subject_content.items():
            content_summary += f"\n📚 {subj}: {info['word_count']} words ({info['chunks_found']} chunks)"
        if not content_summary:
            content_summary = "\n⚠️ No lecture notes found. Plan based on general recommendations."
        
        num_weeks = max(1, min(12, (days_until_exam + 6) // 7))
        
        # Add past paper insights to the prompt
        past_paper_insights = ""
        if past_paper_analysis:
            rankings = get_aggregated_topic_rankings()
            if rankings['topic_rankings']:
                past_paper_insights = "\n📊 Past Paper Topic Rankings (most important first):\n"
                for i, topic in enumerate(rankings['topic_rankings'][:5]):
                    past_paper_insights += f"  {i+1}. {topic['topic']} (Importance: {topic['importance']}, Frequency: {topic['total_frequency']})\n"
        
        system_prompt = f"""You are an expert academic planner. Create a detailed, personalized study plan.
IMPORTANT: Respond ONLY with a valid JSON object. Do NOT wrap it in markdown code blocks.
Do NOT include any text before or after the JSON.

The JSON must have this exact structure:
{{
  "overview": "Brief motivational overview (2-3 sentences)",
  "daily_schedule": "Suggested daily schedule with time blocks",
  "subjects_breakdown": [
    {{
      "subject": "Subject Name",
      "priority": "High",
      "total_hours": 20,
      "topics": ["Topic 1", "Topic 2", "Topic 3"],
      "tips": "Study tips for this subject"
    }}
  ],
  "weekly_plan": [
    {{
      "week": 1,
      "focus": "Main focus for this week",
      "tasks": ["Task 1", "Task 2"]
    }}
  ],
  "revision_strategy": "How to revise effectively",
  "exam_day_tips": "Tips for exam day"
}}

Create exactly {num_weeks} weeks in the weekly_plan.
Each week should have 2-3 specific, actionable tasks.
Make the plan practical and achievable."""
        
        user_prompt = f"""Create a personalized study plan:
📅 Exam Date: {exam_date} ({days_until_exam} days away)
⏰ Study Hours Per Day: {study_hours_per_day} (Total: {total_study_hours}h)
📚 Subjects: {subjects_list}
📄 Notes Found: {content_summary}
Total Content: {total_content_volume} words
Number of Weeks: {num_weeks}
{past_paper_insights}

Create a realistic weekly plan for exactly {num_weeks} weeks."""
        
        response_text = None
        if gemini_api_key:
            for model in GEMINI_MODELS:
                try:
                    response_text = gemini_chat_completion(
                        system_prompt=system_prompt, 
                        contents=[{"role":"user","parts":[{"text":user_prompt}]}], 
                        model=model, 
                        max_tokens=2000, 
                        temperature=0.7
                    )
                    if response_text: break
                except Exception as e:
                    print(f"Gemini study plan error: {e}")
                    continue
        
        if not response_text and groq_api_key:
            for model in GROQ_MODELS:
                try:
                    response_text = groq_chat_completion(
                        messages=[
                            {"role":"system","content":system_prompt},
                            {"role":"user","content":user_prompt}
                        ], 
                        model=model, 
                        max_tokens=2000, 
                        temperature=0.7
                    )
                    if response_text: break
                except Exception as e:
                    print(f"Groq study plan error: {e}")
                    continue
        
        if response_text:
            raw = response_text.strip()
            print(f"Raw study plan response (first 300 chars): {raw[:300]}")
            
            if raw.startswith("```"):
                raw = re.sub(r'^```(?:json)?\s*\n', '', raw)
                raw = re.sub(r'\n```\s*$', '', raw)
            
            try:
                study_plan = json.loads(raw)
                print("✅ Study plan parsed as JSON directly")
            except json.JSONDecodeError:
                json_match = re.search(r'\{[\s\S]*\}', raw)
                if json_match:
                    try:
                        study_plan = json.loads(json_match.group(0))
                        print("✅ Study plan extracted and parsed from response")
                    except json.JSONDecodeError as e:
                        print(f"❌ Failed to extract JSON: {e}")
                        return jsonify({
                            'success': True, 
                            'study_plan': {
                                'overview': raw,
                                'days_until_exam': days_until_exam,
                                'total_study_hours': total_study_hours,
                                'subjects': subjects,
                                'subjects_breakdown': [],
                                'weekly_plan': [],
                                'daily_schedule': '',
                                'revision_strategy': '',
                                'exam_day_tips': '',
                                'raw_response': True
                            }
                        })
                else:
                    return jsonify({
                        'success': True, 
                        'study_plan': {
                            'overview': raw,
                            'days_until_exam': days_until_exam,
                            'total_study_hours': total_study_hours,
                            'subjects': subjects,
                            'subjects_breakdown': [],
                            'weekly_plan': [],
                            'daily_schedule': '',
                            'revision_strategy': '',
                            'exam_day_tips': '',
                            'raw_response': True
                        }
                    })
            
            if isinstance(study_plan.get('overview'), str) and study_plan['overview'].strip().startswith('{'):
                try:
                    nested_json = json.loads(study_plan['overview'])
                    for key in nested_json:
                        if key not in study_plan or not study_plan[key]:
                            study_plan[key] = nested_json[key]
                    if 'overview' in nested_json and not nested_json['overview'].startswith('{'):
                        study_plan['overview'] = nested_json['overview']
                    else:
                        study_plan['overview'] = f"Your {days_until_exam}-day study plan for {', '.join(subjects[:3])} is ready!"
                    print("✅ Nested JSON in overview field was extracted and merged")
                except:
                    pass
            
            study_plan['days_until_exam'] = days_until_exam
            study_plan['total_study_hours'] = total_study_hours
            study_plan['subjects'] = subjects
            
            if 'overview' not in study_plan or not study_plan['overview'] or study_plan['overview'].startswith('{'):
                study_plan['overview'] = f"Your {days_until_exam}-day study plan for {', '.join(subjects[:3])} is ready! Stay focused and you'll succeed."
            
            if 'subjects_breakdown' not in study_plan:
                study_plan['subjects_breakdown'] = []
            if 'weekly_plan' not in study_plan:
                study_plan['weekly_plan'] = []
            if 'daily_schedule' not in study_plan:
                study_plan['daily_schedule'] = ''
            if 'revision_strategy' not in study_plan:
                study_plan['revision_strategy'] = ''
            if 'exam_day_tips' not in study_plan:
                study_plan['exam_day_tips'] = ''
            
            print(f"✅ Study plan finalized: {len(study_plan.get('subjects_breakdown', []))} subjects, {len(study_plan.get('weekly_plan', []))} weeks")
            return jsonify({'success': True, 'study_plan': study_plan})
        
        return jsonify({'error': 'Failed to generate study plan. Please try again.'}), 500
        
    except Exception as e:
        print(f"Study plan error: {traceback.format_exc()}")
        return jsonify({'error': 'An internal error occurred. Please try again.'}), 500

@app.route('/clear-history', methods=['POST'])
def clear_history():
    data = request.json or {}
    clear_session_history(data.get('session_id', 'default'))
    return jsonify({'success': True})

@app.route('/export-chat', methods=['POST'])
def export_chat():
    try:
        data = request.json
        session_id = data.get('session_id', 'default')
        format_type = data.get('format', 'txt')
        messages = get_session_history(session_id, limit=10000)
        if not messages: return jsonify({'error': 'No messages'}), 400
        if format_type == 'txt':
            text = "📝 Chat Export\n" + "="*50 + "\n\n"
            for msg in messages:
                role = "👤 You" if msg['role'] == 'user' else "🤖 AI"
                text += f"{role}: {msg['content']}\n\n"
            return Response(text, mimetype='text/plain', headers={'Content-Disposition': 'attachment;filename=chat_export.txt'})
        elif format_type == 'json':
            return jsonify({'messages': messages, 'exported_at': str(pd.Timestamp.now())})
        else:
            return jsonify({'error': 'Unsupported format'}), 400
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/check-status', methods=['GET'])
def check_status():
    doc_count = 0
    if qdrant_client:
        try: doc_count = qdrant_client.get_collection("university_notes").points_count
        except: pass
    user_count = safe_mongo_count(users_collection)
    return jsonify({
        'status': 'online', 'documents_available': doc_count > 0, 'document_count': doc_count,
        'api_connected': gemini_connected or groq_connected, 'qdrant_connected': qdrant_client is not None,
        'user_system': users_collection is not None, 'total_users': user_count,
        'features': ['persistent_memory', 'streaming', 'follow_up', 'flashcards', 'summaries', 'sentiment', 'suggestions', 'export', 'study_plan', 'past_papers']
    })

@app.route('/admin/documents', methods=['GET'])
@admin_required
def get_documents():
    docs = []
    if qdrant_client:
        try:
            scroll_results = qdrant_client.scroll(collection_name="university_notes", limit=10000, with_payload=True, with_vectors=False)
            file_groups = {}
            for point in scroll_results[0]:
                filename = point.payload.get('filename', '')
                if filename:
                    if filename not in file_groups:
                        file_groups[filename] = {
                            'filename': filename,
                            'file_type': point.payload.get('file_type', filename.split('.')[-1].upper() if '.' in filename else 'WEBPAGE'),
                            'category': point.payload.get('category', 'Root'),
                            'upload_date': point.payload.get('upload_date', 'Unknown'),
                            'doc_id': point.id,
                            'chunks': 0,
                            'text_length': 0
                        }
                    file_groups[filename]['chunks'] += 1
                    file_groups[filename]['text_length'] += len(point.payload.get('text', ''))
            
            # Process quality status for each document
            for info in file_groups.values():
                if info['text_length'] < 100:
                    info['status'] = 'Critical (No text extracted)'
                elif info['chunks'] == 1 and info['file_type'] in ['PDF', 'PPTX'] and info['text_length'] < 400:
                    info['status'] = 'Warning (Low text count)'
                else:
                    info['status'] = 'Healthy (Retrievable)'
            
            docs = list(file_groups.values())
        except Exception as e:
            print(f"Error scrolling documents: {e}")
    return jsonify({'success': True, 'documents': docs})

@app.route('/admin/delete/<doc_id>', methods=['DELETE'])
@admin_required
def delete_document(doc_id):
    try:
        if qdrant_client:
            # Try to retrieve the point to get the filename and delete all associated points
            try:
                res = qdrant_client.retrieve(collection_name="university_notes", ids=[doc_id])
                if res and res[0].payload:
                    filename = res[0].payload.get('filename')
                    if filename:
                        from qdrant_client.http import models as qdrant_models
                        qdrant_client.delete(
                            collection_name="university_notes",
                            points_selector=qdrant_models.Filter(
                                must=[
                                    qdrant_models.FieldCondition(
                                        key="filename",
                                        match=qdrant_models.MatchValue(value=filename)
                                    )
                                ]
                            )
                        )
                        print(f"🗑️ Deleted all chunks for document: {filename}")
                        return jsonify({'success': True})
            except Exception as ex:
                print(f"Failed to delete full document by filename: {ex}. Falling back to single point delete.")
            
            # Fallback to delete by doc_id directly
            qdrant_client.delete(collection_name="university_notes", points_selector=[doc_id])
        return jsonify({'success': True})
    except Exception as e:
        print(f"Delete document error: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/admin/stats')
@admin_required
def get_admin_stats():
    doc_count = 0
    if qdrant_client:
        try: doc_count = qdrant_client.get_collection("university_notes").points_count
        except: pass
    user_count = safe_mongo_count(users_collection)
    past_paper_count = len(past_paper_analysis)
    return jsonify({'success': True, 'total_documents': doc_count, 'total_chunks': doc_count, 'total_users': user_count, 'past_papers_analyzed': past_paper_count})

sync_in_progress = False

@app.route('/admin/sync-website', methods=['POST'])
@admin_required
def admin_sync_website():
    global sync_in_progress
    if sync_in_progress:
        return jsonify({'success': False, 'message': 'Sync already in progress.'}), 400
        
    def perform_sync():
        global sync_in_progress
        sync_in_progress = True
        try:
            from crawler_helper import sync_university_website
            sync_university_website(qdrant_client, embedding_model, max_pages=40)
        finally:
            sync_in_progress = False
            
    import threading
    threading.Thread(target=perform_sync, daemon=True).start()
    return jsonify({'success': True, 'message': 'Sync started in the background.'})

@app.route('/admin/sync-status')
@admin_required
def admin_sync_status():
    global sync_in_progress, qdrant_client
    crawler_count = 0
    if qdrant_client:
        try:
            from qdrant_client.http import models as qdrant_models
            count_res = qdrant_client.count(
                collection_name="university_notes",
                count_filter=qdrant_models.Filter(
                    must=[
                        qdrant_models.FieldCondition(
                            key="source",
                            match=qdrant_models.MatchValue(value="website_crawler")
                        )
                    ]
                )
            )
            crawler_count = count_res.count
        except Exception as e:
            print(f"Error getting sync status count: {e}")
    return jsonify({
        'sync_in_progress': sync_in_progress,
        'crawler_chunks': crawler_count
    })

def run_initial_sync():
    import time
    time.sleep(15)  # Wait for server startup
    global qdrant_client, embedding_model
    if qdrant_client and embedding_model:
        try:
            from qdrant_client.http import models as qdrant_models
            count_res = qdrant_client.count(
                collection_name="university_notes",
                count_filter=qdrant_models.Filter(
                    must=[
                        qdrant_models.FieldCondition(
                            key="source",
                            match=qdrant_models.MatchValue(value="website_crawler")
                        )
                    ]
                )
            )
            if count_res.count == 0:
                print("🔄 [Startup] No university crawled documents found. Starting background crawl...")
                from crawler_helper import sync_university_website
                sync_university_website(qdrant_client, embedding_model, max_pages=35)
            else:
                print(f"ℹ️ [Startup] Found {count_res.count} existing university crawler chunks. Skipping sync.")
        except Exception as e:
            print(f"⚠️ [Startup] Error checking Qdrant for crawler data: {e}")

import threading
threading.Thread(target=run_initial_sync, daemon=True).start()

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 7860))
    doc_count = 0
    if qdrant_client:
        try: doc_count = qdrant_client.get_collection("university_notes").points_count
        except: pass
    user_count = safe_mongo_count(users_collection)
    print("\n" + "="*60)
    print("🚀 SERVER STARTED")
    print(f"📚 Documents: {doc_count}")
    print(f"👥 Users: {user_count}")
    print(f"🧠 Memory: {'MongoDB' if mongo_available else 'SQLite'} + SQLite fallback")
    print(f"📡 Streaming: /chat/stream")
    print(f"🃏 Flashcards: /generate-flashcards")
    print(f"📝 Summaries: /generate-summary")
    print(f"📅 Study Plan: /generate-study-plan")
    print(f"📄 Past Papers: /api/past-papers/rankings")
    print(f"🔍 Web Search: DuckDuckGo + Wikipedia + AnySearch")
    print(f"🌐 http://0.0.0.0:{port}/")
    print("="*60 + "\n")
    app.run(host='0.0.0.0', port=port, debug=False, threaded=True)